from pathlib import Path
from zipfile import ZipFile, ZIP_DEFLATED
import math
import shutil
import tempfile

from lxml import etree
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
VISUALS = ROOT / "visuals"
OUTPUT = ROOT / "final_proposal_satpam.pptx"

WIDE_W = 13.333333
WIDE_H = 7.5

BG = "06142B"
BG_ALT = "081B36"
PANEL = "0B2547"
PANEL_2 = "0D315D"
PANEL_3 = "102E55"
CYAN = "38D9FF"
MINT = "3DDC97"
AMBER = "FFB84D"
RED = "FF4D6D"
BLUE = "1D4ED8"
TEXT = "EAF6FF"
MUTED = "93A8C5"
SOFT = "173A62"
GRID = "123456"
WHITE = "FFFFFF"

TITLE_FONT = "Aptos Display"
BODY_FONT = "Aptos"


def I(value):
    return Inches(value)


def rgb(hex_color):
    hex_color = hex_color.strip("#")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)


def set_line(shape, color=SOFT, width=1.0, dash=None):
    shape.line.color.rgb = rgb(color)
    shape.line.width = Pt(width)
    if dash is not None:
        shape.line.dash_style = dash


def add_text(
    slide,
    x,
    y,
    w,
    h,
    text,
    size=18,
    color=TEXT,
    bold=False,
    font=BODY_FONT,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0.03,
    name=None,
):
    box = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    if name:
        box.name = name
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = I(margin)
    tf.margin_right = I(margin)
    tf.margin_top = I(margin)
    tf.margin_bottom = I(margin)
    tf.vertical_anchor = valign

    lines = text.split("\n")
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(0)
        p.space_before = Pt(0)
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
    return box


def add_panel(slide, x, y, w, h, fill=PANEL, line=SOFT, radius=True, name=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, I(x), I(y), I(w), I(h))
    if name:
        shape.name = name
    set_fill(shape, fill)
    set_line(shape, line, 1.0)
    return shape


def add_rule(slide, x1, y1, x2, y2, color=SOFT, width=1.0):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, I(x1), I(y1), I(x2), I(y2))
    set_line(line, color, width)
    return line


def add_arrow(slide, x1, y1, x2, y2, color=CYAN, width=1.4, head=True):
    line = add_rule(slide, x1, y1, x2, y2, color, width)
    if head:
        head_size = 0.13
        tri = slide.shapes.add_shape(
            MSO_SHAPE.ISOSCELES_TRIANGLE,
            I(x2 - head_size / 2),
            I(y2 - head_size / 2),
            I(head_size),
            I(head_size),
        )
        angle = math.degrees(math.atan2(y2 - y1, x2 - x1))
        tri.rotation = angle + 90
        set_fill(tri, color)
        set_line(tri, color, 0)
    return line


def add_dot(slide, x, y, size=0.11, color=CYAN):
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, I(x - size / 2), I(y - size / 2), I(size), I(size))
    set_fill(dot, color)
    set_line(dot, color, 0)
    return dot


def add_kicker(slide, label):
    marker = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, I(0.58), I(0.44), I(0.12), I(0.12))
    marker.name = "kicker-marker"
    set_fill(marker, CYAN)
    set_line(marker, CYAN, 0)
    add_text(
        slide,
        0.78,
        0.36,
        3.9,
        0.3,
        label.upper(),
        size=10,
        color=CYAN,
        bold=True,
        align=PP_ALIGN.LEFT,
        valign=MSO_ANCHOR.MIDDLE,
        name="kicker-label",
    )


def add_title(slide, kicker, title, subtitle=None):
    add_kicker(slide, kicker)
    add_text(slide, 0.56, 0.78, 8.6, 0.74, title, size=27, color=TEXT, bold=True, font=TITLE_FONT)
    if subtitle:
        add_text(slide, 0.58, 1.43, 8.0, 0.48, subtitle, size=13.5, color=MUTED)


def add_footer(slide, page, source=""):
    add_rule(slide, 0.58, 7.05, 12.75, 7.05, color=GRID, width=0.7)
    add_text(slide, 0.58, 7.11, 8.8, 0.25, source, size=7.8, color=MUTED)
    add_text(slide, 11.8, 7.08, 0.95, 0.28, f"{page:02d} / 10", size=8.5, color=CYAN, bold=True, align=PP_ALIGN.RIGHT)


def add_background(slide, page):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, I(0), I(0), I(WIDE_W), I(WIDE_H))
    set_fill(bg, BG)
    set_line(bg, BG, 0)

    # subtle cyber grid and corner circuitry
    for x in [1.2, 2.4, 3.6, 4.8, 6.0, 7.2, 8.4, 9.6, 10.8, 12.0]:
        add_rule(slide, x, 0.0, x, 7.5, GRID, 0.25)
    for y in [1.2, 2.4, 3.6, 4.8, 6.0]:
        add_rule(slide, 0.0, y, 13.333, y, GRID, 0.25)
    add_rule(slide, 10.4, 0.55, 12.15, 0.55, SOFT, 0.8)
    add_rule(slide, 12.15, 0.55, 12.15, 1.15, SOFT, 0.8)
    add_rule(slide, 11.25, 1.15, 12.15, 1.15, SOFT, 0.8)
    for pos in [(10.4, 0.55), (12.15, 0.55), (12.15, 1.15), (11.25, 1.15)]:
        add_dot(slide, pos[0], pos[1], 0.07, SOFT)
    add_footer(slide, page)


def icon_svg(name, body):
    return f"""<svg xmlns="http://www.w3.org/2000/svg" width="256" height="256" viewBox="0 0 256 256" fill="none">
  <g stroke="#38D9FF" stroke-width="10" stroke-linecap="round" stroke-linejoin="round">
{body}
  </g>
</svg>
"""


def create_mermaid_assets():
    mermaid_theme = """%%{init: {'theme':'dark','themeVariables': {'background':'#06142B','primaryColor':'#0B2547','primaryTextColor':'#EAF6FF','primaryBorderColor':'#38D9FF','lineColor':'#38D9FF','secondaryColor':'#0D315D','tertiaryColor':'#102E55','fontFamily':'Aptos'}}}%%"""

    system_architecture = mermaid_theme + """
flowchart LR
  subgraph Input["Input & Signals"]
    A[Public Reports]
    B[Legal Crawler Findings]
    C[Traffic Simulation]
    D[Transaction Indicators]
    E[Blacklist Dummy]
  end
  subgraph Processing["Ingestion & Processing"]
    F[Validation]
    G[Entity Extraction]
    H[Normalization + Dedup]
    I[Graph Builder]
  end
  subgraph Intelligence["SATPAM Intelligence Core"]
    J[(Neo4j Graph DB)]
    K[A* / BFS / UCS / BDS]
    L[Rule-based Risk Scoring]
    M[Early Warning]
    N[Explainability Engine]
  end
  subgraph Output["Dashboard & Governance"]
    O[Graph Explorer]
    P[Priority Queue]
    Q[Human Verification]
    R[Audit + Export]
  end
  A --> F
  B --> F
  C --> F
  D --> F
  E --> F
  F --> G --> H --> I --> J
  J --> K --> L --> M --> N
  N --> O
  N --> P
  P --> Q --> R
"""

    data_flow = mermaid_theme + """
flowchart LR
  R[Report Text] --> X[Extract URL / WA / Account / APK / Keyword]
  C[Crawler Finding] --> X
  T[Transaction Indicator] --> X
  X --> N[Normalize + Mask Sensitive Fields]
  N --> D[Deduplicate Entity]
  D --> G[Create Node + Relationship]
  G --> S[(Graph Storage)]
  S --> A[Search Path]
  S --> V[Cluster + Centrality]
  A --> O[Evidence Path]
  V --> O
  O --> H[Human Review]
"""

    risk_pipeline = mermaid_theme + """
flowchart LR
  A[Source Signals] --> B[Feature Rules]
  B --> C{Risk Components}
  C --> C1[Report Frequency 30%]
  C --> C2[Blacklist Link 20%]
  C --> C3[Suspicious Transaction 20%]
  C --> C4[Domain / APK Link 15%]
  C --> C5[Fast Fund Movement 15%]
  C1 --> D[Score Cap 100]
  C2 --> D
  C3 --> D
  C4 --> D
  C5 --> D
  D --> E[Risk Level]
  E --> F[Explanation + Review Status]
"""

    (VISUALS / "system_architecture.mmd").write_text(system_architecture.strip() + "\n", encoding="utf-8")
    (VISUALS / "data_flow.mmd").write_text(data_flow.strip() + "\n", encoding="utf-8")
    (VISUALS / "risk_scoring_pipeline.mmd").write_text(risk_pipeline.strip() + "\n", encoding="utf-8")


def create_svg_assets():
    icons = {
        "public_report.svg": icon_svg(
            "public_report",
            """    <path d="M76 42h88l28 28v118a26 26 0 0 1-26 26H76a26 26 0 0 1-26-26V68a26 26 0 0 1 26-26Z"/>
    <path d="M164 42v40h40"/>
    <path d="M86 118h86"/>
    <path d="M86 150h58"/>
    <circle cx="100" cy="188" r="14"/>
    <path d="M124 202c-8-16-40-16-48 0"/>""",
        ),
        "web_crawler.svg": icon_svg(
            "web_crawler",
            """    <circle cx="112" cy="112" r="58"/>
    <path d="M54 112h116"/>
    <path d="M112 54c18 22 18 94 0 116"/>
    <path d="M112 54c-18 22-18 94 0 116"/>
    <path d="M154 154l44 44"/>
    <circle cx="204" cy="204" r="18"/>
    <path d="M54 82h116"/>
    <path d="M54 142h116"/>""",
        ),
        "transaction_monitoring.svg": icon_svg(
            "transaction_monitoring",
            """    <rect x="40" y="62" width="176" height="124" rx="22"/>
    <path d="M40 96h176"/>
    <path d="M68 142h34"/>
    <path d="M128 154l20-32 18 52 20-34 18 18"/>
    <path d="M72 208h112"/>""",
        ),
        "ai_analysis.svg": icon_svg(
            "ai_analysis",
            """    <rect x="70" y="70" width="116" height="116" rx="18"/>
    <path d="M96 116h56"/>
    <path d="M96 146h56"/>
    <circle cx="102" cy="102" r="8"/>
    <circle cx="154" cy="154" r="8"/>
    <path d="M52 96H28"/>
    <path d="M52 128H28"/>
    <path d="M52 160H28"/>
    <path d="M204 96h24"/>
    <path d="M204 128h24"/>
    <path d="M204 160h24"/>
    <path d="M96 52V28"/>
    <path d="M128 52V28"/>
    <path d="M160 52V28"/>
    <path d="M96 228v-24"/>
    <path d="M128 228v-24"/>
    <path d="M160 228v-24"/>""",
        ),
        "risk_score.svg": icon_svg(
            "risk_score",
            """    <path d="M54 162a74 74 0 1 1 148 0"/>
    <path d="M128 162l48-48"/>
    <circle cx="128" cy="162" r="10"/>
    <path d="M64 204h128"/>
    <path d="M76 162h16"/>
    <path d="M128 88v16"/>
    <path d="M180 162h16"/>""",
        ),
        "dashboard.svg": icon_svg(
            "dashboard",
            """    <rect x="36" y="48" width="184" height="152" rx="20"/>
    <path d="M36 86h184"/>
    <rect x="58" y="110" width="58" height="58" rx="10"/>
    <path d="M138 118h54"/>
    <path d="M138 146h42"/>
    <path d="M138 174h54"/>
    <circle cx="70" cy="68" r="4"/>
    <circle cx="88" cy="68" r="4"/>""",
        ),
        "law_enforcement_action.svg": icon_svg(
            "law_enforcement_action",
            """    <path d="M128 34l78 30v58c0 48-28 82-78 102-50-20-78-54-78-102V64l78-30Z"/>
    <path d="M92 130l24 24 58-64"/>
    <path d="M82 210h92"/>""",
        ),
    }
    for filename, content in icons.items():
        (VISUALS / filename).write_text(content, encoding="utf-8")


def draw_icon_png(name, draw_func):
    scale = 4
    size = 256 * scale
    img = Image.new("RGBA", (size, size), (0, 0, 0, 0))
    draw = ImageDraw.Draw(img)
    color = tuple(int(CYAN[i : i + 2], 16) for i in (0, 2, 4)) + (255,)
    accent = tuple(int(MINT[i : i + 2], 16) for i in (0, 2, 4)) + (255,)
    muted = tuple(int(MUTED[i : i + 2], 16) for i in (0, 2, 4)) + (255,)

    def S(v):
        return int(v * scale)

    draw_func(draw, S, color, accent, muted)
    img = img.resize((256, 256), Image.Resampling.LANCZOS)
    img.save(VISUALS / f"{name}.png")


def create_png_assets():
    def report(d, S, c, a, m):
        d.rounded_rectangle([S(50), S(38), S(190), S(218)], radius=S(18), outline=c, width=S(8))
        d.line([S(158), S(38), S(206), S(86), S(158), S(86), S(158), S(38)], fill=c, width=S(8), joint="curve")
        d.line([S(78), S(118), S(164), S(118)], fill=c, width=S(8))
        d.line([S(78), S(148), S(142), S(148)], fill=m, width=S(8))
        d.ellipse([S(84), S(176), S(112), S(204)], outline=a, width=S(8))
        d.arc([S(70), S(194), S(126), S(238)], 205, 335, fill=a, width=S(8))

    def crawler(d, S, c, a, m):
        d.ellipse([S(46), S(46), S(166), S(166)], outline=c, width=S(8))
        d.line([S(46), S(106), S(166), S(106)], fill=c, width=S(6))
        d.arc([S(78), S(46), S(134), S(166)], 70, 290, fill=c, width=S(6))
        d.arc([S(78), S(46), S(134), S(166)], -110, 110, fill=c, width=S(6))
        d.line([S(152), S(152), S(202), S(202)], fill=a, width=S(9))
        d.ellipse([S(186), S(186), S(224), S(224)], outline=a, width=S(8))
        d.line([S(60), S(78), S(152), S(78)], fill=m, width=S(5))
        d.line([S(60), S(136), S(152), S(136)], fill=m, width=S(5))

    def transaction(d, S, c, a, m):
        d.rounded_rectangle([S(36), S(62), S(220), S(184)], radius=S(20), outline=c, width=S(8))
        d.line([S(36), S(96), S(220), S(96)], fill=c, width=S(8))
        d.line([S(66), S(140), S(102), S(140)], fill=m, width=S(8))
        pts = [(S(124), S(154)), (S(146), S(124)), (S(164), S(174)), (S(186), S(138)), (S(210), S(158))]
        d.line(pts, fill=a, width=S(8), joint="curve")
        d.line([S(72), S(208), S(184), S(208)], fill=c, width=S(7))

    def ai(d, S, c, a, m):
        d.rounded_rectangle([S(70), S(70), S(186), S(186)], radius=S(18), outline=c, width=S(8))
        for pos in [96, 128, 160]:
            d.line([S(28), S(pos), S(52), S(pos)], fill=m, width=S(6))
            d.line([S(204), S(pos), S(228), S(pos)], fill=m, width=S(6))
            d.line([S(pos), S(28), S(pos), S(52)], fill=m, width=S(6))
            d.line([S(pos), S(204), S(pos), S(228)], fill=m, width=S(6))
        d.line([S(96), S(116), S(152), S(116)], fill=c, width=S(7))
        d.line([S(96), S(146), S(152), S(146)], fill=c, width=S(7))
        d.ellipse([S(94), S(94), S(112), S(112)], outline=a, width=S(6))
        d.ellipse([S(146), S(146), S(164), S(164)], outline=a, width=S(6))

    def risk(d, S, c, a, m):
        d.arc([S(48), S(54), S(208), S(214)], 190, 340, fill=c, width=S(10))
        d.arc([S(48), S(54), S(208), S(214)], 300, 340, fill=a, width=S(12))
        d.line([S(128), S(162), S(178), S(112)], fill=a, width=S(9))
        d.ellipse([S(118), S(152), S(138), S(172)], fill=a)
        d.line([S(64), S(204), S(192), S(204)], fill=c, width=S(8))
        d.line([S(76), S(162), S(92), S(162)], fill=m, width=S(6))
        d.line([S(128), S(88), S(128), S(104)], fill=m, width=S(6))
        d.line([S(180), S(162), S(196), S(162)], fill=m, width=S(6))

    def dashboard(d, S, c, a, m):
        d.rounded_rectangle([S(36), S(48), S(220), S(200)], radius=S(20), outline=c, width=S(8))
        d.line([S(36), S(86), S(220), S(86)], fill=c, width=S(7))
        d.rounded_rectangle([S(58), S(110), S(116), S(168)], radius=S(10), outline=a, width=S(7))
        for y in [118, 146, 174]:
            d.line([S(138), S(y), S(194), S(y)], fill=m if y == 146 else c, width=S(7))
        d.ellipse([S(66), S(64), S(74), S(72)], fill=a)
        d.ellipse([S(86), S(64), S(94), S(72)], fill=a)

    def action(d, S, c, a, m):
        pts = [(S(128), S(34)), (S(206), S(64)), (S(206), S(122)), (S(188), S(174)), (S(128), S(224)), (S(68), S(174)), (S(50), S(122)), (S(50), S(64))]
        d.line(pts + [pts[0]], fill=c, width=S(8), joint="curve")
        d.line([S(92), S(130), S(116), S(154), S(174), S(90)], fill=a, width=S(11), joint="curve")
        d.line([S(82), S(210), S(174), S(210)], fill=m, width=S(7))

    for name, func in {
        "public_report": report,
        "web_crawler": crawler,
        "transaction_monitoring": transaction,
        "ai_analysis": ai,
        "risk_score": risk,
        "dashboard": dashboard,
        "law_enforcement_action": action,
    }.items():
        draw_icon_png(name, func)


def create_visual_assets():
    VISUALS.mkdir(exist_ok=True)
    create_mermaid_assets()
    create_svg_assets()
    create_png_assets()


def add_icon(slide, name, x, y, w, h=None):
    path = VISUALS / f"{name}.png"
    return slide.shapes.add_picture(str(path), I(x), I(y), width=I(w), height=I(h or w))


def add_chip(slide, x, y, text, color=CYAN, w=None):
    w = w or max(1.0, 0.12 * len(text) + 0.35)
    chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, I(x), I(y), I(w), I(0.34))
    set_fill(chip, BG_ALT)
    set_line(chip, color, 1.0)
    add_text(slide, x + 0.08, y + 0.06, w - 0.16, 0.2, text, size=8.5, color=color, bold=True, align=PP_ALIGN.CENTER)
    return chip


def stat_card(slide, x, y, value, label, accent=CYAN):
    add_panel(slide, x, y, 3.65, 1.0, fill=PANEL, line=accent)
    add_text(slide, x + 0.2, y + 0.15, 3.2, 0.36, value, size=24, color=accent, bold=True, font=TITLE_FONT)
    add_text(slide, x + 0.22, y + 0.58, 3.2, 0.28, label, size=10.2, color=TEXT)


def mini_node(slide, x, y, text, icon, color=CYAN, w=1.35):
    add_panel(slide, x, y, w, 1.0, fill=PANEL, line=color)
    add_icon(slide, icon, x + 0.12, y + 0.16, 0.34)
    add_text(slide, x + 0.5, y + 0.18, w - 0.6, 0.58, text, size=9.5, color=TEXT, bold=True, valign=MSO_ANCHOR.MIDDLE)


def slide_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, 1)
    add_kicker(slide, "Proposal / AI Security Platform")

    add_text(slide, 0.62, 1.05, 5.6, 0.8, "SATPAM", size=50, color=TEXT, bold=True, font=TITLE_FONT)
    add_text(
        slide,
        0.68,
        1.82,
        5.3,
        0.85,
        "Search-based AI Threat Prevention and Mapping",
        size=17,
        color=CYAN,
        bold=True,
    )
    add_text(
        slide,
        0.68,
        2.52,
        5.1,
        0.9,
        "Sistem graph intelligence untuk deteksi, pemetaan, dan prioritisasi risiko ekosistem judol-pinjol ilegal.",
        size=16,
        color=MUTED,
    )

    nodes = [
        (7.2, 1.55, "Report", MINT),
        (9.1, 0.95, "Domain", CYAN),
        (10.9, 1.7, "WA", AMBER),
        (8.2, 3.0, "Rekening", RED),
        (10.4, 3.45, "APK", MINT),
        (6.75, 4.3, "Blacklist", CYAN),
        (9.2, 5.15, "Dashboard", AMBER),
    ]
    links = [(0, 1), (1, 2), (0, 3), (3, 4), (2, 4), (3, 5), (5, 6), (4, 6), (1, 6)]
    for a, b in links:
        add_rule(slide, nodes[a][0], nodes[a][1], nodes[b][0], nodes[b][1], color=SOFT, width=1.1)
    for x, y, label, color in nodes:
        add_dot(slide, x, y, 0.33, color)
        add_text(slide, x - 0.42, y + 0.22, 0.9, 0.22, label, size=8.2, color=TEXT, align=PP_ALIGN.CENTER)

    y = 4.1
    for idx, (label, color) in enumerate(
        [("Detect", CYAN), ("Map", MINT), ("Explain", AMBER), ("Prioritize", RED), ("Verify", CYAN)]
    ):
        add_chip(slide, 0.68 + idx * 1.16, y, label, color, w=1.0)

    add_panel(slide, 0.68, 5.38, 4.95, 0.72, fill=PANEL_2, line=CYAN)
    add_text(slide, 0.9, 5.55, 4.5, 0.26, "Graph Search + Risk Scoring + Human Verification", size=12.2, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 1, "Proposal SATPAM | Prototype berbasis data dummy dan human-in-the-loop")


def slide_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, 2)
    add_title(slide, "Background / Urgency", "Masalahnya bukan satu situs. Ini jaringan lintas kanal.", "Judol-pinjol ilegal bergerak melalui konten, kontak, transaksi, aplikasi, dan laporan korban.")

    stat_card(slide, 0.72, 2.0, "Rp286,84 T", "perputaran dana judol 2025", CYAN)
    stat_card(slide, 4.25, 2.0, "422,1 juta", "transaksi judol tercatat", AMBER)
    stat_card(slide, 7.78, 2.0, "12,3 juta", "orang melakukan deposit", RED)

    add_panel(slide, 0.75, 3.65, 11.85, 1.6, fill=BG_ALT, line=SOFT)
    chain = [
        ("public_report", "Laporan\nmasyarakat", MINT),
        ("web_crawler", "Crawler\nfinding", CYAN),
        ("transaction_monitoring", "Transaksi\nsimulasi", AMBER),
        ("risk_score", "Skor\nrisiko", RED),
        ("dashboard", "Dashboard\nprioritas", MINT),
    ]
    x0 = 1.05
    for i, (icon, label, color) in enumerate(chain):
        x = x0 + i * 2.28
        add_icon(slide, icon, x, 4.0, 0.44)
        add_text(slide, x + 0.52, 3.95, 1.1, 0.5, label, size=9.5, color=TEXT, bold=True)
        if i < len(chain) - 1:
            add_arrow(slide, x + 1.54, 4.22, x + 2.08, 4.22, color=SOFT, width=1.2)

    add_text(slide, 0.86, 5.78, 10.9, 0.46, "Inti gap: sinyal terlihat terpisah, sementara pelaku berpindah domain, nomor, rekening, dan aplikasi.", size=17, color=TEXT, bold=True, font=TITLE_FONT)
    add_footer(slide, 2, "Sumber: PPATK, Catatan Capaian Strategis Tahun 2025, 28 Jan 2026.")


def slide_3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, 3)
    add_title(slide, "Existing System and Gap", "Sistem existing kuat, tetapi relasi lintas ekosistem belum menjadi pusat analisis.", "SATPAM diposisikan sebagai sistem pendukung analisis, bukan pengganti lembaga atau proses resmi.")

    x_left, x_right = 0.74, 6.78
    y0 = 2.0
    add_panel(slide, x_left, y0, 5.65, 4.78, fill=BG_ALT, line=SOFT)
    add_panel(slide, x_right, y0, 5.65, 4.78, fill=BG_ALT, line=CYAN)
    add_text(slide, x_left + 0.22, y0 + 0.18, 5.1, 0.34, "Sistem Existing", size=18, color=MUTED, bold=True, font=TITLE_FONT)
    add_text(slide, x_right + 0.22, y0 + 0.18, 5.1, 0.34, "SATPAM", size=18, color=CYAN, bold=True, font=TITLE_FONT)

    rows = [
        ("Fokus domain masing-masing", "Konten, laporan, transaksi, aplikasi dianalisis sebagai graph terpadu"),
        ("Output: blokir, laporan, atau tindak lanjut", "Output: evidence path, cluster, skor risiko, prioritas"),
        ("Cenderung reaktif setelah entitas ditemukan", "Early warning berbasis pola dan kemiripan jaringan"),
        ("Relasi lintas kanal sulit terlihat", "Domain-WA-rekening-APK-laporan dapat ditelusuri"),
        ("Keputusan rawan overclaim jika otomatis", "Human verification sebelum status final"),
    ]
    for i, (left, right) in enumerate(rows):
        y = y0 + 0.78 + i * 0.74
        add_rule(slide, x_left + 0.2, y - 0.1, x_left + 5.45, y - 0.1, color=GRID, width=0.5)
        add_rule(slide, x_right + 0.2, y - 0.1, x_right + 5.45, y - 0.1, color=GRID, width=0.5)
        add_text(slide, x_left + 0.26, y, 5.0, 0.36, left, size=11.2, color=TEXT)
        add_text(slide, x_right + 0.26, y, 5.0, 0.36, right, size=11.2, color=TEXT)

    add_footer(slide, 3, "Sumber: Komdigi, OJK/IASC, PPATK, dan dokumen proposal SATPAM.")


def slide_4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, 4)
    add_title(slide, "Rumusan Masalah", "Bagaimana sinyal tersebar menjadi jalur risiko yang dapat dijelaskan?", "Rumusan masalah diarahkan pada desain sistem, integrasi data, AI search, scoring, dan explainability.")

    add_panel(slide, 4.38, 2.75, 4.55, 1.18, fill=PANEL_2, line=CYAN)
    add_text(slide, 4.62, 3.04, 4.05, 0.36, "Evidence Path", size=24, color=CYAN, bold=True, font=TITLE_FONT, align=PP_ALIGN.CENTER)
    add_text(slide, 4.72, 3.43, 3.85, 0.2, "laporan -> domain -> WA -> rekening -> APK", size=9.8, color=TEXT, align=PP_ALIGN.CENTER)

    questions = [
        (0.85, 2.05, "1", "Bagaimana merancang graph ekosistem judol-pinjol?", CYAN),
        (8.95, 2.05, "2", "Bagaimana menghubungkan laporan, domain, rekening, APK, dan blacklist?", MINT),
        (0.85, 4.45, "3", "Bagaimana A* dan BFS menemukan jalur risiko?", AMBER),
        (8.95, 4.45, "4", "Bagaimana skor risiko dan prioritas dihitung?", RED),
        (4.1, 5.55, "5", "Bagaimana hasil AI tetap explainable dan diverifikasi manusia?", CYAN),
    ]
    for x, y, num, text, color in questions:
        add_panel(slide, x, y, 3.55, 0.94, fill=BG_ALT, line=color)
        add_dot(slide, x + 0.38, y + 0.47, 0.36, color)
        add_text(slide, x + 0.26, y + 0.35, 0.25, 0.2, num, size=10, color=BG, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, x + 0.68, y + 0.19, 2.85, 0.44, text, size=10.3, color=TEXT)
        add_rule(slide, x + 1.8, y + 0.47, 6.65, 3.34, SOFT, 0.8)

    add_footer(slide, 4, "Rumusan masalah disarikan dari SATPAM_Latar_Belakang_Gap_Inovasi.md.")


def slide_5(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, 5)
    add_title(slide, "Proposed Solution: SATPAM", "SATPAM menyatukan sinyal menjadi peta risiko terverifikasi.", "Alur solusi dibuat sebagai sistem pendukung analisis: indikatif, explainable, dan human-in-the-loop.")

    stages = [
        (0.78, 2.45, "Input\nmulti-sumber", "public_report", MINT),
        (2.72, 2.45, "Entity\nextraction", "ai_analysis", CYAN),
        (4.66, 2.45, "Graph\nintelligence", "web_crawler", AMBER),
        (6.6, 2.45, "A* search\n+ scoring", "risk_score", RED),
        (8.54, 2.45, "Dashboard\nprioritas", "dashboard", CYAN),
        (10.48, 2.45, "Human\nverification", "law_enforcement_action", MINT),
    ]
    for i, (x, y, label, icon, color) in enumerate(stages):
        add_panel(slide, x, y, 1.55, 1.42, fill=PANEL, line=color)
        add_icon(slide, icon, x + 0.47, y + 0.18, 0.48)
        add_text(slide, x + 0.15, y + 0.82, 1.25, 0.42, label, size=10.5, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
        if i < len(stages) - 1:
            add_arrow(slide, x + 1.55, y + 0.71, x + 1.91, y + 0.71, color=SOFT, width=1.2)

    add_panel(slide, 1.08, 4.72, 10.88, 0.78, fill=BG_ALT, line=SOFT)
    add_text(slide, 1.32, 4.94, 10.4, 0.26, "Output utama: risk score, evidence path, cluster, early warning, rekomendasi prioritas, dan status review.", size=13.2, color=TEXT, bold=True, align=PP_ALIGN.CENTER)

    add_text(slide, 1.55, 5.92, 10.1, 0.38, "Batas aman prototype: data dummy, source `simulation_only`, masking data sensitif, dan tidak ada auto-blocking.", size=14, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 5, "Konsep solusi berdasarkan SRS SATPAM bagian scope, architecture, dan human verification.")


def slide_6(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, 6)
    add_title(slide, "System Architecture", "Arsitektur memisahkan ingestion, graph intelligence, scoring, dan review.", "Komponen dirancang modular agar prototype aman, dapat diuji, dan mudah diperluas.")

    lane_y = [1.9, 3.2, 4.55]
    lane_titles = ["Input & Signals", "SATPAM Intelligence Core", "Output & Governance"]
    for y, title in zip(lane_y, lane_titles):
        add_panel(slide, 0.72, y, 11.85, 1.05, fill=BG_ALT, line=SOFT)
        add_text(slide, 0.92, y + 0.12, 1.35, 0.22, title, size=9.2, color=CYAN, bold=True)

    # Lane 1
    mini_node(slide, 2.18, 1.93, "Report\nForm/API", "public_report", MINT, 1.45)
    mini_node(slide, 3.98, 1.93, "Crawler\nFinding", "web_crawler", CYAN, 1.45)
    mini_node(slide, 5.78, 1.93, "Traffic\nSimulasi", "transaction_monitoring", AMBER, 1.45)
    mini_node(slide, 7.58, 1.93, "Transaksi\nSimulasi", "transaction_monitoring", RED, 1.45)
    mini_node(slide, 9.38, 1.93, "Blacklist\nDummy", "risk_score", CYAN, 1.45)

    # Lane 2
    core = [
        (1.48, "Validation"),
        (2.82, "Extraction"),
        (4.16, "Normalize\n+ Dedup"),
        (5.5, "Graph\nBuilder"),
        (6.84, "Neo4j\nGraph DB"),
        (8.18, "A*/BFS\nUCS/BDS"),
        (9.52, "Risk\nScoring"),
        (10.86, "Explain\n+ Alert"),
    ]
    for i, (x, label) in enumerate(core):
        add_panel(slide, x, 3.34, 1.05, 0.62, fill=PANEL_2 if i in [4, 5] else PANEL, line=CYAN if i in [4, 5] else SOFT)
        add_text(slide, x + 0.06, 3.47, 0.93, 0.22, label, size=8.8, color=TEXT, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        if i < len(core) - 1:
            add_arrow(slide, x + 1.05, 3.65, x + 1.28, 3.65, SOFT, 1.0)

    # Lane 3
    outputs = [
        (2.0, "Graph\nExplorer", "dashboard", CYAN),
        (4.1, "Priority\nQueue", "risk_score", AMBER),
        (6.2, "Verification\nCase", "law_enforcement_action", MINT),
        (8.3, "Audit\nLog", "ai_analysis", CYAN),
        (10.4, "Export\nReport", "public_report", MINT),
    ]
    for x, label, icon, color in outputs:
        mini_node(slide, x, 4.58, label, icon, color, 1.5)
    add_arrow(slide, 10.95, 2.95, 10.95, 3.19, SOFT, 1.0)
    add_arrow(slide, 10.95, 4.0, 10.95, 4.55, SOFT, 1.0)

    add_footer(slide, 6, "Mermaid source: visuals/system_architecture.mmd | Diagram PPTX digambar sebagai shape editable.")


def slide_7(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, 7)
    add_title(slide, "Method / Workflow", "Pipeline metode berjalan dari data dummy sampai rekomendasi prioritas.", "Metode menggabungkan graph database, AI search, rule-based risk scoring, dan review manusia.")

    steps = [
        ("1", "Input\nData", "laporan, crawler,\ntransaksi simulasi", MINT),
        ("2", "Extract", "URL, domain,\nWA, rekening, APK", CYAN),
        ("3", "Normalize", "masking,\ndedup entity", AMBER),
        ("4", "Build\nGraph", "node +\nrelationship", CYAN),
        ("5", "Search", "A*, BFS,\nUCS, BDS", RED),
        ("6", "Score", "rule-based\nrisk score", AMBER),
        ("7", "Explain\nReview", "path bukti +\nhuman approval", MINT),
    ]
    x0 = 0.75
    for i, (num, title, desc, color) in enumerate(steps):
        x = x0 + i * 1.78
        add_dot(slide, x + 0.52, 2.65, 0.44, color)
        add_text(slide, x + 0.39, 2.52, 0.25, 0.2, num, size=9.8, color=BG, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, x + 0.08, 3.08, 1.0, 0.46, title, size=14, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, x - 0.03, 3.68, 1.2, 0.48, desc, size=8.6, color=MUTED, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            add_arrow(slide, x + 0.78, 2.65, x + 1.35, 2.65, SOFT, 1.1)

    add_panel(slide, 1.05, 5.0, 3.3, 0.88, fill=PANEL, line=MINT)
    add_text(slide, 1.25, 5.18, 2.9, 0.24, "Safety by design", size=15, color=MINT, bold=True, font=TITLE_FONT, align=PP_ALIGN.CENTER)
    add_text(slide, 1.2, 5.5, 2.95, 0.18, "data dummy, no live interception", size=8.8, color=TEXT, align=PP_ALIGN.CENTER)

    add_panel(slide, 4.95, 5.0, 3.3, 0.88, fill=PANEL, line=CYAN)
    add_text(slide, 5.15, 5.18, 2.9, 0.24, "Explainable output", size=15, color=CYAN, bold=True, font=TITLE_FONT, align=PP_ALIGN.CENTER)
    add_text(slide, 5.1, 5.5, 2.95, 0.18, "rule aktif + evidence path", size=8.8, color=TEXT, align=PP_ALIGN.CENTER)

    add_panel(slide, 8.85, 5.0, 3.3, 0.88, fill=PANEL, line=AMBER)
    add_text(slide, 9.05, 5.18, 2.9, 0.24, "Human-in-the-loop", size=15, color=AMBER, bold=True, font=TITLE_FONT, align=PP_ALIGN.CENTER)
    add_text(slide, 9.0, 5.5, 2.95, 0.18, "candidate dulu, approval kemudian", size=8.8, color=TEXT, align=PP_ALIGN.CENTER)

    add_footer(slide, 7, "Mermaid source: visuals/data_flow.mmd dan visuals/risk_scoring_pipeline.mmd.")


def slide_8(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, 8)
    add_title(slide, "Innovation / Novelty", "Novelty SATPAM: network-based intelligence yang explainable.", "Perbedaan utama bukan hanya deteksi, tetapi pemetaan hubungan dan prioritas berbasis graph risk.")

    add_panel(slide, 0.72, 1.9, 8.85, 4.9, fill=BG_ALT, line=SOFT)
    add_text(slide, 0.95, 2.12, 1.5, 0.26, "Gap", size=12, color=MUTED, bold=True)
    headers = ["Graph", "Search", "Score", "Explain", "Verify"]
    for i, h in enumerate(headers):
        add_text(slide, 3.02 + i * 1.1, 2.12, 0.85, 0.26, h, size=10, color=CYAN, bold=True, align=PP_ALIGN.CENTER)

    rows = [
        ("Data tersebar", [1, 1, 0, 1, 0]),
        ("Sistem reaktif", [1, 1, 1, 0, 1]),
        ("Identitas mudah berganti", [1, 1, 1, 1, 0]),
        ("Judol-pinjol terpisah", [1, 1, 1, 1, 0]),
        ("Output sulit dijelaskan", [1, 1, 1, 1, 1]),
    ]
    for r, (gap, vals) in enumerate(rows):
        y = 2.68 + r * 0.72
        add_rule(slide, 0.92, y - 0.15, 9.28, y - 0.15, GRID, 0.5)
        add_text(slide, 0.95, y, 2.05, 0.24, gap, size=10.2, color=TEXT)
        for i, val in enumerate(vals):
            color = MINT if val else SOFT
            add_dot(slide, 3.42 + i * 1.1, y + 0.11, 0.2 if val else 0.12, color)

    add_panel(slide, 9.92, 1.9, 2.55, 4.9, fill=PANEL_2, line=CYAN)
    add_text(slide, 10.18, 2.25, 2.05, 0.58, "Inovasi inti", size=22, color=CYAN, bold=True, font=TITLE_FONT, align=PP_ALIGN.CENTER)
    bullets = [
        "Graph intelligence lintas entitas",
        "A* risk path untuk prioritas",
        "Judol-pinjol linkage detection",
        "Rule score yang bisa dijelaskan",
        "Human verification by design",
    ]
    for i, bullet in enumerate(bullets):
        add_dot(slide, 10.18, 3.25 + i * 0.5, 0.12, [CYAN, MINT, AMBER, RED, CYAN][i])
        add_text(slide, 10.35, 3.16 + i * 0.5, 1.9, 0.26, bullet, size=9.5, color=TEXT)

    add_footer(slide, 8, "Novelty disarikan dari SATPAM_Latar_Belakang_Gap_Inovasi.md bagian Gap dan Inovasi.")


def slide_9(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, 9)
    add_title(slide, "Benefits and Implementation Roadmap", "Manfaat analitik jelas, roadmap prototype tetap realistis.", "SATPAM memberi prioritas dan jalur bukti tanpa mengambil alih keputusan manusia.")

    add_panel(slide, 0.72, 1.88, 4.2, 4.95, fill=BG_ALT, line=SOFT)
    add_text(slide, 0.98, 2.16, 3.7, 0.34, "Manfaat", size=21, color=CYAN, bold=True, font=TITLE_FONT)
    benefits = [
        ("Prioritas kasus", "analis fokus pada risiko tertinggi", CYAN),
        ("Evidence path", "alasan risiko dapat ditelusuri", MINT),
        ("Cluster jaringan", "melihat node pusat dan relasi", AMBER),
        ("Governance", "candidate, review, audit log", RED),
    ]
    for i, (title, desc, color) in enumerate(benefits):
        y = 2.86 + i * 0.82
        add_dot(slide, 1.03, y + 0.16, 0.22, color)
        add_text(slide, 1.28, y, 3.3, 0.22, title, size=12.5, color=TEXT, bold=True)
        add_text(slide, 1.28, y + 0.27, 3.3, 0.2, desc, size=8.8, color=MUTED)

    add_panel(slide, 5.25, 1.88, 7.22, 4.95, fill=BG_ALT, line=CYAN)
    add_text(slide, 5.52, 2.16, 4.2, 0.34, "Roadmap Prototype", size=21, color=CYAN, bold=True, font=TITLE_FONT)
    roadmap = [
        ("F1", "Foundation", "schema, Docker, seed", 1.0, MINT),
        ("F2", "Data + Graph", "extract, normalize, Neo4j", 1.0, CYAN),
        ("F3", "Search + Scoring", "A*, BFS, rules, alert", 1.0, AMBER),
        ("F4", "Dashboard + Review", "graph explorer, cases", 1.0, RED),
        ("F5", "Testing + Demo", "pytest, UI smoke, script", 0.7, MINT),
    ]
    base_x, base_y = 5.72, 3.05
    for i, (phase, title, desc, width, color) in enumerate(roadmap):
        x = base_x + i * 1.25
        add_dot(slide, x, base_y, 0.28, color)
        if i < len(roadmap) - 1:
            add_rule(slide, x + 0.18, base_y, x + 1.06, base_y, SOFT, 1.4)
        add_text(slide, x - 0.25, base_y + 0.36, 0.5, 0.22, phase, size=10, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, x - 0.45, base_y + 0.68, 0.9, 0.32, title, size=8.6, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, x - 0.55, base_y + 1.12, 1.1, 0.34, desc, size=7.4, color=MUTED, align=PP_ALIGN.CENTER)

    add_panel(slide, 5.82, 5.72, 5.9, 0.48, fill=PANEL, line=AMBER)
    add_text(slide, 6.1, 5.84, 5.35, 0.18, "Target MVP: 1.000 node, 5.000 relationship, dashboard < 3 detik pada dataset prototype.", size=9.5, color=TEXT, bold=True, align=PP_ALIGN.CENTER)

    add_footer(slide, 9, "Roadmap berdasarkan SRS SATPAM fase implementasi prototype.")


def slide_10(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, 10)
    add_kicker(slide, "Closing and References")

    add_text(slide, 0.78, 1.05, 7.9, 0.86, "SATPAM membaca ancaman sebagai jaringan risiko.", size=34, color=TEXT, bold=True, font=TITLE_FONT)
    add_text(slide, 0.82, 2.02, 7.15, 0.66, "Bukan alat vonis otomatis, melainkan decision-support untuk mendeteksi, memetakan, menjelaskan, memprioritaskan, dan memverifikasi.", size=16, color=MUTED)

    words = [("Detect", CYAN), ("Map", MINT), ("Explain", AMBER), ("Prioritize", RED), ("Verify", CYAN)]
    for i, (word, color) in enumerate(words):
        add_chip(slide, 0.86 + i * 1.35, 3.05, word, color, w=1.15)

    add_panel(slide, 0.86, 4.08, 6.85, 1.28, fill=PANEL_2, line=CYAN)
    add_text(slide, 1.12, 4.34, 6.35, 0.5, "Thesis: graph search + risk scoring + human verification dapat membuat analisis judol-pinjol lebih terhubung, transparan, dan prioritas.", size=15.5, color=TEXT, bold=True, font=TITLE_FONT, align=PP_ALIGN.CENTER)

    add_panel(slide, 8.35, 1.1, 3.95, 5.45, fill=BG_ALT, line=SOFT)
    add_text(slide, 8.62, 1.38, 3.3, 0.36, "Referensi utama", size=18, color=CYAN, bold=True, font=TITLE_FONT)
    refs = [
        "[1] PPATK, Catatan Capaian Strategis 2025, 2026.",
        "[2] Komdigi, laporan pemblokiran rekening judol, 2025.",
        "[3] OJK, Indonesia Anti-Scam Centre, 2026.",
        "[4] OJK, Waspada Pinjaman Online Ilegal.",
        "[5]-[8] Neo4j Documentation: graph concepts, path finding, BFS, centrality.",
        "[9]-[11] SRS dan dokumen proposal SATPAM lokal.",
    ]
    for i, ref in enumerate(refs):
        add_text(slide, 8.62, 1.95 + i * 0.62, 3.2, 0.36, ref, size=8.6, color=TEXT)

    add_footer(slide, 10, "Daftar lengkap tersedia di references.md.")


def add_fade_transitions(pptx_path):
    ns = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
    tmp_dir = Path(tempfile.mkdtemp(prefix="satpam_pptx_"))
    try:
        with ZipFile(pptx_path, "r") as zin:
            zin.extractall(tmp_dir)

        slides_dir = tmp_dir / "ppt" / "slides"
        for slide_xml in sorted(slides_dir.glob("slide*.xml")):
            tree = etree.parse(str(slide_xml))
            root = tree.getroot()
            for old in root.xpath("./p:transition", namespaces=ns):
                root.remove(old)
            transition = etree.Element("{http://schemas.openxmlformats.org/presentationml/2006/main}transition")
            transition.set("spd", "slow")
            fade = etree.SubElement(transition, "{http://schemas.openxmlformats.org/presentationml/2006/main}fade")
            c_sld = root.find("p:cSld", namespaces=ns)
            insert_at = 1
            if c_sld is not None:
                insert_at = list(root).index(c_sld) + 1
            root.insert(insert_at, transition)
            tree.write(str(slide_xml), encoding="UTF-8", xml_declaration=True, standalone=True)

        backup = pptx_path.with_suffix(".pptx.bak")
        if backup.exists():
            backup.unlink()
        pptx_path.rename(backup)
        with ZipFile(pptx_path, "w", ZIP_DEFLATED) as zout:
            for file in tmp_dir.rglob("*"):
                if file.is_file():
                    zout.write(file, file.relative_to(tmp_dir).as_posix())
        backup.unlink()
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)


def build_pptx():
    create_visual_assets()

    prs = Presentation()
    prs.slide_width = I(WIDE_W)
    prs.slide_height = I(WIDE_H)
    blank = prs.slide_layouts[6]
    # Remove default empty slide if any are added by templates: Presentation starts with none.
    _ = blank

    slide_1(prs)
    slide_2(prs)
    slide_3(prs)
    slide_4(prs)
    slide_5(prs)
    slide_6(prs)
    slide_7(prs)
    slide_8(prs)
    slide_9(prs)
    slide_10(prs)

    prs.save(OUTPUT)
    add_fade_transitions(OUTPUT)
    verify_pptx(OUTPUT, expected_slides=10)
    print(f"Generated {OUTPUT.name} with 10 slides.")


def verify_pptx(path, expected_slides):
    if not path.exists() or path.stat().st_size == 0:
        raise RuntimeError("PPTX was not created or is empty.")
    with ZipFile(path, "r") as zf:
        slide_parts = [n for n in zf.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml")]
        media_parts = [n for n in zf.namelist() if n.startswith("ppt/media/")]
        if len(slide_parts) != expected_slides:
            raise RuntimeError(f"Expected {expected_slides} slides, found {len(slide_parts)}.")
        for media in media_parts:
            if zf.getinfo(media).file_size == 0:
                raise RuntimeError(f"Empty media file detected: {media}")


if __name__ == "__main__":
    build_pptx()
